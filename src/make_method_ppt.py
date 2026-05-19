"""분석 방법 메모 PPT — Hampel filter + 물리 한계 인용 정리.

특징:
  - 설명 그림 포함 (figures/fig_hampel_*.png, fig_robust_*.png, fig_physical_*.png)
  - 글은 단계 번호로 체계적 구성, 그림과 좌우 분할
  - 간략하지만 핵심은 다 포함
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT = '/Users/gimhanseo/Desktop/공프/자동분석폴더/분석방법_메모.pptx'
FIG_DIR = '/Users/gimhanseo/Desktop/공프/자동분석폴더/figures'

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

BG = RGBColor(0x0A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
TEXT = RGBColor(0xE8, 0xEE, 0xF6)
MUTED = RGBColor(0x8A, 0x9B, 0xB0)
CARD = RGBColor(0x12, 0x2A, 0x45)


def bg(s):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG


def tb(s, text, x, y, w, h, sz=14, c=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(x, y, w, h); tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = bold
    p.font.name = 'Calibri'


def numbered_steps(s, items, x, y, w, h, sz=12, c=TEXT):
    box = s.shapes.add_textbox(x, y, w, h); tf = box.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.name = 'Calibri'
        p.space_after = Pt(7)


def card(s, x, y, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = CARD; sh.line.fill.background()


def line(s, x, y, w):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(3))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()


# ─── Slide 1: cover ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, '분석 방법 메모', Inches(0.7), Inches(2.4), Inches(12), Inches(1.2),
   sz=44, c=TEXT, bold=True)
line(s, Inches(0.7), Inches(3.7), Inches(2))
tb(s, 'Hampel filter · 물리 한계 설정 · 인용 출처',
   Inches(0.7), Inches(4.0), Inches(12), Inches(0.6), sz=18, c=TEXT)
tb(s, '(외부 자료 인용용 간략 정리)',
   Inches(0.7), Inches(4.7), Inches(12), Inches(0.5), sz=13, c=MUTED)

# ─── Slide 2: Hampel filter — 메커니즘 그림 + 단계 설명 ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, '① Hampel Filter — 어떻게 작동하는가',
   Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
   sz=24, c=TEXT, bold=True)
line(s, Inches(0.6), Inches(1.15), Inches(1.5))

# Left: figure
img = os.path.join(FIG_DIR, 'fig_hampel_mechanism.png')
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(0.3), Inches(1.4),
                          Inches(7.6), Inches(5.6))

# Right: 단계
card(s, Inches(8.1), Inches(1.4), Inches(5.0), Inches(5.6))
tb(s, '판정 단계', Inches(8.3), Inches(1.55), Inches(4.7), Inches(0.4),
   sz=15, c=ACCENT, bold=True)
numbered_steps(s, [
    'Step 1.  같은 웨이퍼 내 8방향 이웃 다이 수집',
    '',
    'Step 2.  이웃들의 중간값 m, MAD s 계산',
    '         (m = median, s = median(|x − m|))',
    '',
    'Step 3.  대상 다이 값 x와 비교:',
    '         |x − m| > k · s   →   outlier',
    '         (k = 3, 보통 3-sigma 자리)',
    '',
    'Step 4.  이웃이 < 3개면 검사 skip',
    '         (가장자리 다이 보호)',
], Inches(8.3), Inches(2.0), Inches(4.7), Inches(5.0), sz=11)

# ─── Slide 3: 왜 3-sigma 대신 Hampel — 비교 그림 ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, '② 왜 3-sigma 대신 Hampel을 쓰는가',
   Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
   sz=24, c=TEXT, bold=True)
line(s, Inches(0.6), Inches(1.15), Inches(1.5))

img = os.path.join(FIG_DIR, 'fig_robust_vs_3sigma.png')
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(0.4), Inches(1.4),
                          Inches(12.5), Inches(4.5))

card(s, Inches(0.4), Inches(6.0), Inches(12.5), Inches(1.3))
tb(s, '핵심 정리', Inches(0.6), Inches(6.08), Inches(12), Inches(0.4),
   sz=14, c=ACCENT, bold=True)
numbered_steps(s, [
    '①  3-sigma는 평균·σ 기반 → 이상치가 평균·σ를 끌어올려서 자기 자신을 숨김',
    '②  중간값·MAD는 이상치에 끌려가지 않음 (robust, breakdown 50%)',
    '③  표본 작은 우리 데이터(웨이퍼당 14다이)에선 강건한 통계가 필수',
], Inches(0.6), Inches(6.45), Inches(12), Inches(0.85), sz=11)

# ─── Slide 4: 물리 한계 — 그림 + 표 ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, '③ 물리 한계 (Physical Bounds)',
   Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
   sz=24, c=TEXT, bold=True)
line(s, Inches(0.6), Inches(1.15), Inches(1.5))

img = os.path.join(FIG_DIR, 'fig_physical_bounds.png')
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(0.4), Inches(1.35),
                          Inches(7.0), Inches(4.0))

card(s, Inches(7.6), Inches(1.35), Inches(5.5), Inches(4.0))
tb(s, '디바이스 사양', Inches(7.8), Inches(1.5),
   Inches(5.2), Inches(0.4), sz=14, c=ACCENT, bold=True)
numbered_steps(s, [
    '·  Standard TE-mode Si depletion MZM',
    '·  Single-arm DC bias drive',
    '·  MMI splitter (분배비 ±0.3 dB)',
    '·  Single-stage (cascade 없음)',
    '·  No polarization filter',
    '·  No push-pull, no advanced ER tricks',
    '',
    '→ baseline 디바이스 → 한계도 baseline',
], Inches(7.8), Inches(1.95), Inches(5.2), Inches(3.4), sz=11)

# 표
table_data = [
    ['지표',     '범위',         '근거 요약'],
    ['ER (dB)',  '[0, 45]',      '정의상 ≥0 / TE Si depletion + MMI 한계'],
    ['IL (dB)',  '[−20, 0]',     '패시브 → ≤0 / 운용 한계 마진'],
    ['V_π (V)',  '[5, 80]',      'V_π·L 1~3 V·cm × 짧은 phase shifter / 외삽 한계'],
]
n_rows, n_cols = len(table_data), 3
tbl = s.shapes.add_table(n_rows, n_cols, Inches(0.4), Inches(5.55),
                          Inches(12.5), Inches(1.6)).table
for i, w in enumerate([Inches(1.3), Inches(1.5), Inches(9.7)]):
    tbl.columns[i].width = w
for r, row in enumerate(table_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c); cell.text = val
        p = cell.text_frame.paragraphs[0]; p.font.name = 'Calibri'
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x06, 0x5A, 0x82)
            p.font.color.rgb = TEXT; p.font.bold = True; p.font.size = Pt(12)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0D, 0x1F, 0x36) if r % 2 else RGBColor(0x11, 0x27, 0x40)
            p.font.color.rgb = TEXT; p.font.size = Pt(11)

# ─── Slide 5: 인용 출처 ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, '④ 인용 출처', Inches(0.6), Inches(0.4),
   Inches(12), Inches(0.7), sz=24, c=TEXT, bold=True)
line(s, Inches(0.6), Inches(1.15), Inches(1.5))

card(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.3))
tb(s, '[Si MZM 성능 범위 — ER, IL, V_π]',
   Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
   sz=14, c=ACCENT, bold=True)
numbered_steps(s, [
    '· Patel et al., "Design, analysis, and transmission system performance',
    '         of a 41 GHz silicon photonic modulator",  Optics Express 23, 14263 (2015)',
    '         → standard TE Si MZM의 ER 35–45 dB / IL 4–8 dB 보고',
    '',
    '· Witzens, "High-Speed Silicon Photonics Modulators",',
    '         Proceedings of the IEEE 106, 2158 (2018)',
    '         DOI: 10.1109/JPROC.2018.2877636   → 종합 review',
], Inches(0.7), Inches(2.0), Inches(12), Inches(1.7), sz=11)

card(s, Inches(0.5), Inches(3.85), Inches(12.3), Inches(1.6))
tb(s, '[Plasma Dispersion 원전 — V_π·L 산출 근거]',
   Inches(0.7), Inches(3.97), Inches(12), Inches(0.4),
   sz=14, c=ACCENT, bold=True)
numbered_steps(s, [
    '· Soref & Bennett, "Electrooptical effects in silicon",',
    '         IEEE Journal of Quantum Electronics 23, 123 (1987)',
    '         → Si plasma dispersion 효과의 정량 관계식',
], Inches(0.7), Inches(4.4), Inches(12), Inches(1.0), sz=11)

card(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.6))
tb(s, '[Hampel Filter 원전]',
   Inches(0.7), Inches(5.67), Inches(12), Inches(0.4),
   sz=14, c=ACCENT, bold=True)
numbered_steps(s, [
    '· Hampel, F.R., "The influence curve and its role in robust estimation",',
    '         Journal of the American Statistical Association 69, 383 (1974)',
    '· Pearson, R.K., "Mining Imperfect Data", SIAM (2005) — 실무 가이드',
], Inches(0.7), Inches(6.1), Inches(12), Inches(1.0), sz=11)

# ─── Slide 6: 발표용 인용 템플릿 ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, '⑤ 발표용 인용 템플릿 (한 문장)',
   Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
   sz=24, c=TEXT, bold=True)
line(s, Inches(0.6), Inches(1.15), Inches(1.5))

card(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.6))
tb(s, '물리 한계 정당화', Inches(0.7), Inches(1.62),
   Inches(12), Inches(0.4), sz=14, c=ACCENT, bold=True)
tb(s,
   '"본 디바이스는 standard TE-mode Si depletion MZM (single-arm drive, '
   'MMI splitter, single-stage, 특수 ER 향상 기법 미적용)으로, '
   '동일 사양의 보고 범위 (Patel et al., Opt. Express 2015; '
   'Witzens, Proc. IEEE 2018)와 plasma dispersion 표준값 '
   '(Soref & Bennett, IEEE JQE 1987)을 참고하여 '
   'ER ∈ [0, 45] dB, IL ∈ [−20, 0] dB, V_π ∈ [5, 80] V로 설정."',
   Inches(0.7), Inches(2.05), Inches(12), Inches(2.0), sz=12, c=TEXT)

card(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.6))
tb(s, 'Hampel filter 정당화', Inches(0.7), Inches(4.42),
   Inches(12), Inches(0.4), sz=14, c=ACCENT, bold=True)
tb(s,
   '"이웃 다이 기반 outlier 검출은 Hampel(JASA 1974)의 robust estimation '
   '방법에 기반하며, 같은 웨이퍼 내 8방향 이웃의 중간값/MAD를 기준으로 '
   '|x − median| > 3·MAD인 다이를 outlier로 표기. '
   '3-sigma는 표본이 작고(n=14) 이상치가 평균·σ를 왜곡하는 문제로 채택하지 않음."',
   Inches(0.7), Inches(4.85), Inches(12), Inches(2.0), sz=12, c=TEXT)

prs.save(OUT)
print(f'Saved: {OUT}')
